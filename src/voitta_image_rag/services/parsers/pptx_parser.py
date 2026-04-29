"""PPTX parser via python-pptx.

Each slide becomes a section in the markdown output. Text-frame contents are
concatenated; pictures are extracted with position = start of the slide's
section.
"""

from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import ClassVar

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base import BaseParser, ExtractedImage, ParserResult

logger = logging.getLogger(__name__)


class PptxParser(BaseParser):
    extensions: ClassVar[list[str]] = [".pptx"]

    def parse(self, file_path: Path) -> ParserResult:
        try:
            prs = Presentation(str(file_path))
        except Exception as e:
            return ParserResult.failure(f"pptx open failed: {e}")

        parts: list[str] = []
        images: list[ExtractedImage] = []
        cursor = 0

        for slide_index, slide in enumerate(prs.slides, start=1):
            slide_start = cursor
            heading = f"# Slide {slide_index}"
            text_lines: list[str] = [heading]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs)
                        if line:
                            text_lines.append(line)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        blob: bytes = shape.image.blob
                        ct: str = shape.image.content_type
                    except Exception as e:  # pragma: no cover — defensive
                        logger.warning("pptx image read failed: %s", e)
                        continue
                    width, height = _dimensions(blob)
                    images.append(
                        ExtractedImage(
                            bytes=blob,
                            mime=ct,
                            position=slide_start,
                            page=slide_index,
                            width=width,
                            height=height,
                        )
                    )

            section = "\n".join(text_lines)
            parts.append(section)
            cursor += len(section) + 2  # joiner

        content = "\n\n".join(parts)
        return ParserResult(content=content, images=images)


def _dimensions(data: bytes) -> tuple[int | None, int | None]:
    try:
        from PIL import Image as PILImage

        with PILImage.open(io.BytesIO(data)) as img:
            return img.size
    except Exception:
        return (None, None)
